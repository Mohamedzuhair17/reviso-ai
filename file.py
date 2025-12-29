import io
import os
import time
import streamlit as st
import PyPDF2
import google.generativeai as genai
from pptx import Presentation
from dotenv import load_dotenv
from google.api_core.exceptions import ResourceExhausted, PermissionDenied

# ================= LOAD ENV =================
load_dotenv()

API_KEY = os.getenv("GEMINI_API_KEY")
if not API_KEY:
    st.error("GEMINI_API_KEY missing. Add it to Streamlit Secrets.")
    st.stop()

genai.configure(api_key=API_KEY)

# ================= FORCE SAFE MODEL =================
model = genai.GenerativeModel("models/gemini-1.5-flash")

# ================= PAGE CONFIG =================
st.set_page_config(
    page_title="Reviso — The AI School Assistant",
    layout="wide"
)

# ================= SAFE GENERATION =================
def safe_generate(prompt, retries=2):
    for _ in range(retries):
        try:
            return model.generate_content(prompt).text
        except ResourceExhausted:
            time.sleep(5)
        except PermissionDenied:
            return (
                "❌ Permission denied.\n\n"
                "Check API enablement and billing."
            )
    return "⚠️ AI is busy. Try again later."

# ================= SESSION =================
if "mode" not in st.session_state:
    st.session_state.mode = None
if "output" not in st.session_state:
    st.session_state.output = None
if "ppt" not in st.session_state:
    st.session_state.ppt = None
if "last_run" not in st.session_state:
    st.session_state.last_run = 0

# ================= UI =================
st.title("Reviso")
st.caption("The AI School Assistant — From notes to understanding")

uploaded = st.file_uploader("Upload notes (PDF)", type=["pdf"])

if uploaded:
    reader = PyPDF2.PdfReader(io.BytesIO(uploaded.read()))
    notes = " ".join(page.extract_text() or "" for page in reader.pages)

    if not notes.strip():
        st.error("No readable text found in PDF.")
        st.stop()

    # ---------- MODES ----------
    c1, c2, c3, c4 = st.columns(4)

    with c1:
        if st.button("🧠 Understand"):
            st.session_state.mode = "summary"
            st.session_state.output = None
            st.session_state.ppt = None

    with c2:
        if st.button("📄 Revise"):
            st.session_state.mode = "cheat"
            st.session_state.output = None
            st.session_state.ppt = None

    with c3:
        if st.button("📊 Present"):
            st.session_state.mode = "ppt"
            st.session_state.output = None
            st.session_state.ppt = None

    with c4:
        if st.button("❓ Ask"):
            st.session_state.mode = "qa"
            st.session_state.output = None
            st.session_state.ppt = None

    # ---------- COOLDOWN ----------
    if time.time() - st.session_state.last_run < 8:
        st.warning("Please wait a few seconds before generating again.")
        st.stop()

    st.session_state.last_run = time.time()

    # ---------- WORKING ----------
    if st.session_state.mode == "summary" and st.session_state.output is None:
        with st.spinner("Explaining your notes..."):
            st.session_state.output = safe_generate(
                f"""
                Create a deep, exam-oriented explanation.
                Use headings, definitions, and examples.

                Notes:
                {notes[:6000]}
                """
            )

    if st.session_state.mode == "cheat" and st.session_state.output is None:
        with st.spinner("Preparing revision sheet..."):
            st.session_state.output = safe_generate(
                f"""
                Create a ONE-PAGE exam cheat sheet.
                Bullet points only.

                Notes:
                {notes[:6000]}
                """
            )

    if st.session_state.mode == "ppt" and st.session_state.ppt is None:
        with st.spinner("Preparing presentation..."):
            outline = safe_generate(
                f"""
                Create a professional PPT outline.
                Max 8 slides. Short bullets.

                Notes:
                {notes[:5000]}
                """
            )

            prs = Presentation()
            title = prs.slides.add_slide(prs.slide_layouts[0])
            title.shapes.title.text = "Study Notes"
            title.placeholders[1].text = "Prepared using Reviso"

            for block in outline.split("\n\n"):
                lines = block.split("\n")
                if len(lines) < 2:
                    continue
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = lines[0][:60]
                slide.placeholders[1].text = "\n".join(lines[1:])

            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)
            st.session_state.ppt = buf

    # ---------- QUESTIONING FEATURE ----------
    if st.session_state.mode == "qa":
        st.markdown("### Ask a question from your notes")
        question = st.text_input("Type your question")

        if question:
            with st.spinner("Finding answer..."):
                answer = safe_generate(
                    f"""
                    Answer the question strictly using the notes below.
                    If not found, say so clearly.

                    Notes:
                    {notes[:6000]}

                    Question:
                    {question}
                    """
                )
                st.session_state.output = answer

    # ---------- OUTPUT ----------
    if st.session_state.output:
        st.markdown("### Output")
        st.write(st.session_state.output)

    if st.session_state.ppt:
        st.download_button(
            "⬇️ Download PPT",
            st.session_state.ppt,
            file_name="reviso_presentation.pptx"
        )

else:
    st.info("Upload your notes to begin.")

